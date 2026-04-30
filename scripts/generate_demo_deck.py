#!/usr/bin/env python3
"""Generate the Standards Control Plane demo deck (PowerPoint).

Reads the architecture diagrams from `docs/diagrams/` (rendered from
mermaid via `docs/diagrams/render.sh`) and assembles a walk-through deck
mirroring `docs/decks/SCP-overview-2026-04-30.md`.

Usage:
    python3 scripts/generate_demo_deck.py
    # writes docs/demo-guide.pptx

Re-run any time the markdown overview changes — the deck is generated,
not hand-edited. To re-render the architecture diagrams from their
mermaid source, run `docs/diagrams/render.sh` first.

Convention follows mapp-size-allocation/scripts/generate_demo_deck.py
and mapp-returns-intelligence/scripts/generate_demo_deck.py.
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parent.parent
DIAGRAMS = REPO / "docs" / "diagrams"
OUT = REPO / "docs" / "demo-guide.pptx"

# 16:9 widescreen — 13.333 × 7.5 inches
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette — picked to match the rendered mermaid diagrams
INK = RGBColor(0x10, 0x14, 0x29)
INK_MUTED = RGBColor(0x4B, 0x55, 0x6B)
PAPER = RGBColor(0xFA, 0xFA, 0xF7)
SCP_BLUE = RGBColor(0x1E, 0x40, 0xAF)
SCP_INDIGO = RGBColor(0x10, 0x14, 0x40)
ACC_RUST = RGBColor(0x9A, 0x34, 0x12)
REPO_GREEN = RGBColor(0x14, 0x53, 0x2D)
ACCENT_TEAL = RGBColor(0x0E, 0x76, 0x90)
ACCENT_AMBER = RGBColor(0xB4, 0x53, 0x09)
ACCENT_RED = RGBColor(0x99, 0x1B, 0x1B)
ACCENT_VIOLET = RGBColor(0x6D, 0x28, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xD1, 0xD5, 0xDB)
STRIPE = RGBColor(0xF3, 0xF4, 0xF6)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def add_blank_slide(prs: Presentation, fill: RGBColor = PAPER):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = fill
    bg.shadow.inherit = False
    return slide


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_eyebrow(slide, text, color=SCP_BLUE):
    add_text(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
             text, size=12, bold=True, color=color)


def add_title(slide, text, color=INK):
    add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
             text, size=36, bold=True, color=color)


def add_subtitle(slide, text, color=INK_MUTED):
    add_text(slide, Inches(0.6), Inches(1.65), Inches(12), Inches(0.5),
             text, size=16, color=color)


def add_image_fit(slide, path: Path, left, top, width, height,
                  *, bg: RGBColor | None = None):
    """Fit image into the box preserving aspect ratio, centred. If `bg`
    is provided, paint a backing rectangle first (useful for diagrams
    that were rendered with a transparent background)."""
    if bg is not None:
        backing = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        backing.fill.solid()
        backing.fill.fore_color.rgb = bg
        backing.line.fill.background()
        backing.shadow.inherit = False

    if not path.exists():
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        ph.fill.solid()
        ph.fill.fore_color.rgb = STRIPE
        ph.line.color.rgb = INK_MUTED
        add_text(slide, left, top + height // 2 - Inches(0.2), width, Inches(0.4),
                 f"[missing: {path.name}]", size=10, color=INK_MUTED, align=PP_ALIGN.CENTER)
        return

    with Image.open(path) as im:
        iw, ih = im.size
    box_ratio = width / height
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        # image is wider than box — fit by width
        out_w = width
        out_h = Emu(int(width * ih / iw))
        out_left = left
        out_top = top + (height - out_h) // 2
    else:
        out_h = height
        out_w = Emu(int(height * iw / ih))
        out_top = top
        out_left = left + (width - out_w) // 2
    slide.shapes.add_picture(str(path), out_left, out_top, width=out_w, height=out_h)


def add_bullets(slide, left, top, width, height, items, *, size=14, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"


def add_chip(slide, left, top, text, *, fill=SCP_BLUE, fg=WHITE, width=Inches(1.6)):
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.32))
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill
    chip.line.fill.background()
    chip.shadow.inherit = False
    tf = chip.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = fg
    run.font.name = "Calibri"


def add_table(slide, left, top, headers, rows, col_widths, *, row_h=Inches(0.55), header_color=INK_MUTED):
    """Estate-standard simple table: header row in muted small caps, then
    striped body rows, with a thin rule under the header."""
    x = left
    for h, w in zip(headers, col_widths):
        add_text(slide, x, top, w, Inches(0.4), h, size=11, bold=True, color=header_color)
        x += w
    rule = slide.shapes.add_connector(1, left, top + Inches(0.4),
                                      left + sum(col_widths, Inches(0)), top + Inches(0.4))
    rule.line.color.rgb = RULE

    for i, row in enumerate(rows):
        row_top = top + Inches(0.55) + row_h * i
        if i % 2 == 0:
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, row_top - Inches(0.05),
                                            sum(col_widths, Inches(0)), row_h)
            stripe.fill.solid(); stripe.fill.fore_color.rgb = STRIPE
            stripe.line.fill.background()
        x = left
        for cell, w in zip(row, col_widths):
            cell_text = cell["text"] if isinstance(cell, dict) else str(cell)
            cell_size = cell.get("size", 12) if isinstance(cell, dict) else 12
            cell_color = cell.get("color", INK) if isinstance(cell, dict) else INK
            cell_bold = cell.get("bold", False) if isinstance(cell, dict) else False
            add_text(slide, x + Inches(0.1), row_top, w - Inches(0.2), row_h,
                     cell_text, size=cell_size, color=cell_color, bold=cell_bold)
            x += w


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------


def slide_title(prs):
    s = add_blank_slide(prs, fill=SCP_INDIGO)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.6), Inches(0.4), Inches(2.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = SCP_BLUE; bar.line.fill.background()
    add_text(s, Inches(1.4), Inches(2.5), Inches(11), Inches(0.5),
             "STANDARDS CONTROL PLANE", size=14, bold=True, color=SCP_BLUE)
    add_text(s, Inches(1.4), Inches(3.0), Inches(11), Inches(1.5),
             "Demo overview", size=66, bold=True, color=WHITE)
    add_text(s, Inches(1.4), Inches(4.4), Inches(11), Inches(0.6),
             "What it does · how it's enforced · how it evolves",
             size=22, color=RGBColor(0xC0, 0xC8, 0xE8))
    add_text(s, Inches(1.4), Inches(6.5), Inches(11), Inches(0.4),
             "James Brooke · 2026-04-30 · walk-through draft",
             size=12, color=RGBColor(0x90, 0x9C, 0xC4))


def slide_problem(prs):
    s = add_blank_slide(prs)
    add_eyebrow(s, "WHY THIS EXISTS")
    add_title(s, "Three pains we hit before SCP")
    add_subtitle(s, "Multiple repos · multiple humans · multiple AI agents — each drifting in parallel.")
    cols = [
        ("DRIFT",
         "schema fragmentation",
         "Each repo's services.yml, waiver shape, vendoring conventions diverged. Reviews argued style, not substance.",
         ACCENT_AMBER),
        ("AGENTS INVENTING STANDARDS",
         "plausible ≠ ours",
         "Codex / Sonnet / Opus would pick a plausible pattern but not necessarily our pattern. Review caught it; the wasted dispatch didn't.",
         ACCENT_VIOLET),
        ("NO ENFORCEMENT FLOOR",
         "decisions ≠ enforcement",
         "DECISIONS.md was aspirational — nothing forced a downstream PR to honour it. Governance was retrospective.",
         ACCENT_RED),
    ]
    col_w = Inches(3.9)
    gap = Inches(0.3)
    start = Inches(0.6)
    top = Inches(2.6)
    for i, (head, sub, body, color) in enumerate(cols):
        left = start + (col_w + gap) * i
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_w, Inches(4.2))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = color; card.line.width = Pt(1.5)
        accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_w, Inches(0.18))
        accent.fill.solid(); accent.fill.fore_color.rgb = color; accent.line.fill.background()
        add_text(s, left + Inches(0.3), top + Inches(0.5), col_w - Inches(0.6), Inches(0.6),
                 head, size=22, bold=True, color=color)
        add_text(s, left + Inches(0.3), top + Inches(1.2), col_w - Inches(0.6), Inches(0.4),
                 sub, size=14, bold=True, color=INK_MUTED)
        add_text(s, left + Inches(0.3), top + Inches(1.8), col_w - Inches(0.6), Inches(2.2),
                 body, size=14, color=INK)
    add_text(s, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
             "The bet: make standards machine-readable + machine-enforceable. Drift collapses; review handles novel substance.",
             size=12, bold=True, color=INK_MUTED)


def slide_what_scp_is(prs):
    s = add_blank_slide(prs)
    add_eyebrow(s, "WHAT SCP IS")
    add_title(s, "Three sentences")
    add_subtitle(s, "")
    pitch = [
        "SCP is a deterministic policy plane.",
        "It exposes \"what's allowed\" in two flavours: a consult API that agents query before writing code,",
        "and a federation gate that blocks merges after writing code if the result violates a rule.",
    ]
    for i, line in enumerate(pitch):
        add_text(s, Inches(0.6), Inches(2.3) + Inches(0.7) * i, Inches(12), Inches(0.65),
                 line, size=24, bold=(i == 0), color=INK if i == 0 else INK_MUTED)
    # Owns / does NOT own table
    top = Inches(4.8)
    add_text(s, Inches(0.6), top, Inches(6.2), Inches(0.4),
             "WHAT SCP OWNS", size=11, bold=True, color=SCP_BLUE)
    add_text(s, Inches(7.1), top, Inches(6.2), Inches(0.4),
             "WHAT SCP DOES NOT OWN", size=11, bold=True, color=INK_MUTED)
    rule_y = top + Inches(0.4)
    rule = s.shapes.add_connector(1, Inches(0.6), rule_y, Inches(13.0 - 0.6), rule_y)
    rule.line.color.rgb = RULE

    owns = [
        "Source of truth for governance + architecture standards",
        "Findings store, waivers, scorecards",
        "Consult interface (humans + agents)",
        "Federation primitive (reusable workflow downstream pins)",
    ]
    not_owns = [
        "Multi-agent orchestration → ACC",
        "Service runtime / auth fabric → Control Tower",
        "Documentation retrieval → mapp-doc-agent",
        "The estate itself — SCP plugs in, doesn't replace",
    ]
    add_bullets(s, Inches(0.6), Inches(5.3), Inches(6.2), Inches(2.0), owns, size=13)
    add_bullets(s, Inches(7.1), Inches(5.3), Inches(6.2), Inches(2.0), not_owns, size=13, color=INK_MUTED)


def slide_section(prs, *, label, body, fill=SCP_INDIGO, accent=SCP_BLUE):
    s = add_blank_slide(prs, fill=fill)
    add_text(s, Inches(0.6), Inches(2.8), Inches(12), Inches(0.6),
             label, size=18, bold=True, color=accent)
    add_text(s, Inches(0.6), Inches(3.3), Inches(12), Inches(2),
             body, size=44, bold=True, color=WHITE)


def slide_diagram(prs, *, eyebrow, title, subtitle, image, caption,
                  color=SCP_BLUE, image_box=None):
    """Standard diagram slide: eyebrow + title + subtitle + centred
    image (fitted, aspect-preserved) + caption underneath."""
    s = add_blank_slide(prs)
    add_eyebrow(s, eyebrow, color=color)
    add_title(s, title)
    if subtitle:
        add_subtitle(s, subtitle)
    if image_box is None:
        image_box = (Inches(1.0), Inches(2.4), Inches(11.3), Inches(4.5))
    left, top, width, height = image_box
    add_image_fit(s, image, left, top, width, height)
    add_text(s, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4),
             caption, size=11, color=INK_MUTED, align=PP_ALIGN.CENTER)


def slide_diagram_with_bullets(prs, *, eyebrow, title, image, bullets,
                               color=SCP_BLUE):
    """Diagram on the left, bullet list on the right."""
    s = add_blank_slide(prs)
    add_eyebrow(s, eyebrow, color=color)
    add_title(s, title)
    add_image_fit(s, image, Inches(0.6), Inches(2.0), Inches(7.6), Inches(5.0))
    add_text(s, Inches(8.6), Inches(2.0), Inches(4.4), Inches(0.4),
             "WHY", size=11, bold=True, color=color)
    add_bullets(s, Inches(8.6), Inches(2.4), Inches(4.4), Inches(4.6), bullets, size=13)


def slide_three_rules(prs):
    s = add_blank_slide(prs)
    add_eyebrow(s, "WHAT v1.0.0 ACTUALLY CHECKS")
    add_title(s, "Three rules — deliberately small")
    add_subtitle(s, "The primitive matters more than the rule library. Rules 4–N follow via RFC in v1.1+.")
    headers = ["RULE", "WHAT IT CHECKS", "WHY THIS RULE FIRST"]
    col_widths = [Inches(2.0), Inches(5.3), Inches(5.0)]
    rows = [
        [
            {"text": "SCP-R-001", "bold": True, "color": SCP_BLUE, "size": 14},
            {"text": "services.yml root-shape conforms to SVC-003 auth modes", "size": 12},
            {"text": "Auth-contract drift is highest blast radius. WP-SCP-019 froze the schema; this rule enforces it.", "size": 12, "color": INK_MUTED},
        ],
        [
            {"text": "SCP-R-002", "bold": True, "color": SCP_BLUE, "size": 14},
            {"text": "waivers.json entry schema (approved_by, created_at, expires_at, rule_id ∨ finding_id)", "size": 12},
            {"text": "Waivers are the safety valve. If their shape drifts, the safety valve becomes opaque.", "size": 12, "color": INK_MUTED},
        ],
        [
            {"text": "SCP-R-003", "bold": True, "color": SCP_BLUE, "size": 14},
            {"text": "Vendoring-manifest marker present in package.json / pyproject.toml / go.mod when those change", "size": 12},
            {"text": "Marks every adopter as having consciously adopted the SVC-003 contract.", "size": 12, "color": INK_MUTED},
        ],
    ]
    add_table(s, Inches(0.6), Inches(2.5), headers, rows, col_widths, row_h=Inches(0.95))
    add_text(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4),
             "Each rule: structured deny payload · allow + deny fixtures · opa test coverage ≥ 90% · waiver-aware via policies/scp_common.rego",
             size=11, color=INK_MUTED, align=PP_ALIGN.CENTER)


def slide_threshold_a(prs):
    s = add_blank_slide(prs)
    add_eyebrow(s, "THE FINISH LINE", color=ACCENT_TEAL)
    add_title(s, "Threshold A — what \"actually useful\" means")
    add_subtitle(s, "SCP gates itself on its own main · 3 rules enforcing · conflict-gate green · v1.0.0 cut.")
    # Two-column status
    done = [
        ("020A   plan + D-022/D-023", True),
        ("020B   reusable workflow (PR #36)", True),
        ("020B.1 workflow-selftest harness (PR #38)", True),
        ("020B.2 scripts/scp-policy-check (PR #41)", True),
        ("020C   3 starter Rego rules (PR #49)", True),
        ("020C.1 waiver-aware + conflict-gate (PR #52)", True),
        ("020J   tag-protection + signed-commits (PR #53)", True),
    ]
    todo = [
        ("020K   adopter onboarding + ADOPT-005", False),
        ("020D1  self-dogfood wrapper (advisory)", False),
        ("020H.1 v1.0.0-rc.1 + observability emit", False),
        ("020E.a pre-protection canary", False),
        ("🛑 USER-GATE-A0 advisory→required signoff", False),
        ("020D2  required-status-check + cut v1.0.0", False),
        ("🛑 USER-GATE-A Threshold A signoff", False),
    ]
    top = Inches(2.4)
    add_text(s, Inches(0.6), top, Inches(6.2), Inches(0.4),
             "✅ DONE", size=11, bold=True, color=ACCENT_TEAL)
    add_text(s, Inches(7.1), top, Inches(6.2), Inches(0.4),
             "⏳ AHEAD", size=11, bold=True, color=ACCENT_AMBER)
    rule_y = top + Inches(0.4)
    rule = s.shapes.add_connector(1, Inches(0.6), rule_y, Inches(13.0 - 0.6), rule_y)
    rule.line.color.rgb = RULE
    for i, (item, _) in enumerate(done):
        add_text(s, Inches(0.6), top + Inches(0.55) + Inches(0.4) * i,
                 Inches(6.2), Inches(0.4), item, size=12,
                 color=INK, font="Consolas")
    for i, (item, _) in enumerate(todo):
        add_text(s, Inches(7.1), top + Inches(0.55) + Inches(0.4) * i,
                 Inches(6.2), Inches(0.4), item, size=12,
                 color=INK_MUTED, font="Consolas")


def slide_state_today(prs):
    s = add_blank_slide(prs)
    add_eyebrow(s, "WHERE WE ARE TODAY")
    add_title(s, "Programmes + active PRs")
    add_subtitle(s, "Snapshot: 2026-04-30 morning, post-020J apply.")
    headers = ["PROGRAMME", "STATE", "TRACK"]
    col_widths = [Inches(7.0), Inches(3.4), Inches(2.0)]
    rows = [
        [{"text": "WP-SCP-019 — Service Auth Contract", "bold": True, "size": 13},
         {"text": "✅ closed 2026-04-20", "size": 12, "color": ACCENT_TEAL},
         {"text": "—", "size": 12, "color": INK_MUTED}],
        [{"text": "WP-SCP-020 — Federation Primitive", "bold": True, "size": 13},
         {"text": "🔄 in flight (~60% slices)", "size": 12, "color": ACCENT_AMBER},
         {"text": "Track 1", "size": 12, "color": INK_MUTED}],
        [{"text": "WP-SCP-021 — MCP Server", "bold": True, "size": 13},
         {"text": "✅ landed 2026-04-29", "size": 12, "color": ACCENT_TEAL},
         {"text": "Track 2", "size": 12, "color": INK_MUTED}],
        [{"text": "WP-SCP-022 — Implementation Programme", "bold": True, "size": 13},
         {"text": "🔄 orchestrator for 020 + 021", "size": 12, "color": ACCENT_AMBER},
         {"text": "Meta", "size": 12, "color": INK_MUTED}],
    ]
    add_table(s, Inches(0.6), Inches(2.5), headers, rows, col_widths, row_h=Inches(0.65))
    # Live infra
    add_text(s, Inches(0.6), Inches(5.4), Inches(12), Inches(0.4),
             "LIVE INFRA (today)", size=11, bold=True, color=SCP_BLUE)
    rule_y = Inches(5.8)
    rule = s.shapes.add_connector(1, Inches(0.6), rule_y, Inches(13.0 - 0.6), rule_y)
    rule.line.color.rgb = RULE
    infra = [
        "main has required_signatures: true (020J landed this morning)",
        "Tag-protection ruleset on v* — no force-push, no deletion, no re-pointing",
        "3 Rego rules under policies/, reusable workflow at .github/workflows/policy-check.yml",
        "MCP server at src/standards_control_plane/mcp/, Ed25519 signing key under .scp/keys/",
    ]
    add_bullets(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(1.4), infra, size=12)


def slide_benefits(prs):
    s = add_blank_slide(prs)
    add_eyebrow(s, "WHY THIS IS WORTH DOING")
    add_title(s, "Benefits — three audiences")
    add_subtitle(s, "Estate · operator · future agents and humans.")
    cols = [
        ("FOR THE ESTATE",
         [
             "Drift collapses — every repo held to the same schema",
             "Agents stop inventing — every dispatch starts with \"what's allowed?\"",
             "Review handles novel substance, not boilerplate",
             "Governance becomes reproducible — policies/ is a git tree",
         ],
         SCP_BLUE),
        ("FOR THE OPERATOR",
         [
             "Single source of truth, NOT single point of failure",
             "Renovate cascades pin-bumps automatically",
             "Bus-factor-1 acknowledged explicitly (quarterly review cadence)",
             "Conflict-gate is structural insurance",
         ],
         ACCENT_TEAL),
        ("FOR FUTURE AGENTS",
         [
             "Standards evolve, never deprecate suddenly (RFC + 1-release window)",
             "Receipt-bound consults are auditable forensic trails",
             "Machine-readable rules = pre-code conformance",
             "MCP server is advisory; gate is authoritative — belt + braces",
         ],
         ACCENT_VIOLET),
    ]
    col_w = Inches(3.9)
    gap = Inches(0.3)
    start = Inches(0.6)
    top = Inches(2.6)
    for i, (head, items, color) in enumerate(cols):
        left = start + (col_w + gap) * i
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_w, Inches(4.5))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = color; card.line.width = Pt(1.5)
        accent_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_w, Inches(0.18))
        accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = color
        accent_bar.line.fill.background()
        add_text(s, left + Inches(0.3), top + Inches(0.45), col_w - Inches(0.6), Inches(0.5),
                 head, size=15, bold=True, color=color)
        add_bullets(s, left + Inches(0.3), top + Inches(1.1), col_w - Inches(0.6), Inches(3.2),
                    items, size=12)


def slide_evolution_timeline(prs):
    s = add_blank_slide(prs)
    add_eyebrow(s, "HOW IT EVOLVES")
    add_title(s, "Phased rollout — a timeline of trust")
    add_subtitle(s, "Each phase is one mechanism. Build the floor first; optimise on top.")
    headers = ["PHASE", "WHAT CHANGES", "RISK"]
    col_widths = [Inches(3.0), Inches(7.5), Inches(1.9)]
    rows = [
        [{"text": "0  advisory audit", "size": 12, "bold": True},
         {"text": "Post-merge audit. Findings logged, no enforcement.", "size": 11, "color": INK_MUTED},
         {"text": "ZERO", "size": 11, "bold": True, "color": ACCENT_TEAL}],
        [{"text": "1  auth contract", "size": 12, "bold": True},
         {"text": "WP-SCP-019. SVC-003 schema frozen. Three modes + one deprecating.", "size": 11, "color": INK_MUTED},
         {"text": "LOW", "size": 11, "bold": True, "color": ACCENT_TEAL}],
        [{"text": "2  federation primitive", "size": 12, "bold": True},
         {"text": "WP-SCP-020. Reusable workflow + Rego bundle + Renovate cascade. SCP self-pilot.", "size": 11, "color": INK_MUTED},
         {"text": "MEDIUM", "size": 11, "bold": True, "color": ACCENT_AMBER}],
        [{"text": "3  MCP server", "size": 12, "bold": True},
         {"text": "WP-SCP-021. Agents consult before writing code. Track 2 closed.", "size": 11, "color": INK_MUTED},
         {"text": "LOW", "size": 11, "bold": True, "color": ACCENT_TEAL}],
        [{"text": "4  self-gate", "size": 12, "bold": True},
         {"text": "020D1 → 020D2. SCP gates its own main. Advisory → required. v1.0.0 cut.", "size": 11, "color": INK_MUTED},
         {"text": "HIGH", "size": 11, "bold": True, "color": ACCENT_RED}],
        [{"text": "5  multi-adopter", "size": 12, "bold": True},
         {"text": "FLA pilot first, then estate cascade. Renovate cascades pin-bumps.", "size": 11, "color": INK_MUTED},
         {"text": "MEDIUM", "size": 11, "bold": True, "color": ACCENT_AMBER}],
        [{"text": "6  scorecards", "size": 12, "bold": True},
         {"text": "WP-SCP-023. Cross-repo aggregation, proposal queue, D-NNN indexing.", "size": 11, "color": INK_MUTED},
         {"text": "LOW", "size": 11, "bold": True, "color": ACCENT_TEAL}],
    ]
    add_table(s, Inches(0.6), Inches(2.4), headers, rows, col_widths, row_h=Inches(0.6))


def slide_reference(prs):
    s = add_blank_slide(prs, fill=SCP_INDIGO)
    add_text(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.5),
             "REFERENCE", size=14, bold=True, color=SCP_BLUE)
    add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.9),
             "Where to dig deeper", size=42, bold=True, color=WHITE)
    refs = [
        ("docs/decks/SCP-overview-2026-04-30.md",
         "the markdown source of truth for this deck"),
        ("docs/plans/WP-SCP-020-policy-federation-primitive.md",
         "federation primitive design, slice-by-slice"),
        ("docs/plans/WP-SCP-021-scp-as-mcp-server.md",
         "MCP server design + receipt model"),
        ("docs/plans/WP-SCP-022-implementation-programme-plan.md",
         "how 020 + 021 are sequenced + gate-by-gate roadmap"),
        ("docs/adoption/ADOPT-001-project-onboarding.md",
         "the single onboarding brief for adopter repos"),
        ("docs/integrations/conflict-gate.md",
         "conflict-gate adapter doctrine"),
        ("docs/security/branch-protection.md",
         "020J configured-state docs + verify + revert"),
        ("docs/DECISIONS.md",
         "every amending D-NNN decision row"),
        ("policies/README.md",
         "rule template + 7-step contributor checklist"),
        ("STATUS.md",
         "at-a-glance Threshold A progress + active PRs"),
    ]
    top = Inches(2.5)
    for i, (path, desc) in enumerate(refs):
        row_top = top + Inches(0.42) * i
        add_text(s, Inches(0.6), row_top, Inches(5.5), Inches(0.4),
                 path, size=12, color=RGBColor(0xC0, 0xC8, 0xE8), font="Consolas")
        add_text(s, Inches(6.3), row_top, Inches(7.0), Inches(0.4),
                 desc, size=12, color=RGBColor(0x90, 0x9C, 0xC4))


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Title
    slide_title(prs)

    # 2. The problem
    slide_problem(prs)

    # 3. What SCP is
    slide_what_scp_is(prs)

    # 4. Section: ARCHITECTURE
    slide_section(prs,
                  label="ARCHITECTURE",
                  body="How the pieces fit together.",
                  fill=SCP_INDIGO,
                  accent=SCP_BLUE)

    # 5. Three-layer architecture
    slide_diagram(prs,
                  eyebrow="ARCHITECTURE · THREE LAYERS",
                  title="Policy plane · orchestration plane · adopter plane",
                  subtitle="SCP plugs into the estate; it doesn't try to be the estate.",
                  image=DIAGRAMS / "diagram-01.png",
                  caption="SCP sits above ACC (orchestration) and adopter repos (execution). Two protocols out: MCP for agents · reusable-workflow for PRs.",
                  color=SCP_BLUE)

    # 6. Two operating modes
    slide_diagram(prs,
                  eyebrow="OPERATING MODES",
                  title="Consult before code · audit after code",
                  subtitle="Same registry, same rules — one source of truth.",
                  image=DIAGRAMS / "diagram-02.png",
                  caption="Pre-code consult closes the dispatch-cycle waste; post-code audit catches drift at PR time. Belt + braces.",
                  color=SCP_BLUE)

    # 7. Section: ENFORCEMENT
    slide_section(prs,
                  label="ENFORCEMENT",
                  body="The federation primitive — how a PR\nactually gets blocked.",
                  fill=SCP_INDIGO,
                  accent=ACCENT_AMBER)

    # 8. Federation primitive stack
    slide_diagram(prs,
                  eyebrow="ENFORCEMENT · FEDERATION PRIMITIVE",
                  title="Reusable workflow + Rego + branch protection",
                  subtitle="Adopter wrapper is ~20 lines. Pinned by SHA. Renovate cascades the bump.",
                  image=DIAGRAMS / "diagram-03.png",
                  caption="OPA + Conftest pinned by SHA256. Caller's GITHUB_TOKEN is the ceiling. Break-glass needs three gates — fails closed.",
                  color=ACCENT_AMBER)

    # 9. The three rules
    slide_three_rules(prs)

    # 10. Conflict-gate
    slide_diagram(prs,
                  eyebrow="ENFORCEMENT · CONFLICT-GATE",
                  title="Two engines must agree, or merge is blocked",
                  subtitle="Architectural insurance against \"the rule looks enforced but silently doesn't fire.\"",
                  image=DIAGRAMS / "diagram-04.png",
                  caption="Rego = fast PR gate · Python = deep audit. Disagreement → SCP-E005 → amending D-NNN in a separate PR.",
                  color=ACCENT_AMBER)

    # 11. Section: AGENTS
    slide_section(prs,
                  label="AGENTS",
                  body="The MCP surface — closing the loop\nbefore code is written.",
                  fill=SCP_INDIGO,
                  accent=ACCENT_VIOLET)

    # 12. MCP server flow
    slide_diagram_with_bullets(prs,
                               eyebrow="AGENTS · MCP SERVER",
                               title="Pre-code consult with signed receipts",
                               image=DIAGRAMS / "diagram-05.png",
                               bullets=[
                                   "Agent calls scp.consult_rules(scope) before writing code",
                                   "Server returns rules + Ed25519-signed receipt (TTL ≤ 2h)",
                                   "Receipt binds (rule-set hash, scope, timestamp)",
                                   "Pre-flight check via scp.propose(diff, receipt) is advisory",
                                   "Reusable workflow at PR time stays authoritative",
                                   "Track 2 closed 2026-04-29 (USER-GATE-C signed)",
                               ],
                               color=ACCENT_VIOLET)

    # 13. Section: ESTATE
    slide_section(prs,
                  label="ESTATE",
                  body="How SCP plugs into the rest\nof Brokapps.",
                  fill=SCP_INDIGO,
                  accent=REPO_GREEN)

    # 14. Estate integration map
    slide_diagram(prs,
                  eyebrow="ESTATE · INTEGRATION MAP",
                  title="Policy plane · orchestration · runtime · adopters",
                  subtitle="Boundary discipline: SCP doesn't try to be ACC + CT + doc-agent + every repo.",
                  image=DIAGRAMS / "diagram-06.png",
                  caption="ACC consults pre-code · adopters call the gate at PR time · CT + doc-agent integrate post-Threshold A.",
                  color=REPO_GREEN)

    # 15. Section: EVOLUTION
    slide_section(prs,
                  label="EVOLUTION",
                  body="How it changes over time —\nphases of trust.",
                  fill=SCP_INDIGO,
                  accent=ACCENT_TEAL)

    # 16. Phased rollout (table)
    slide_evolution_timeline(prs)

    # 17. Threshold A
    slide_threshold_a(prs)

    # 18. State today
    slide_state_today(prs)

    # 19. Post-v1.0.0 evolution arrow
    slide_diagram(prs,
                  eyebrow="POST-v1.0.0 ROADMAP",
                  title="v1.x → v2.x — what comes after self-gating",
                  subtitle="Each version adds capability without breaking adopters. Pin-bump via Renovate is the upgrade path.",
                  image=DIAGRAMS / "diagram-08.png",
                  caption="v1.1: FLA pilot · v1.2: estate cascade · v2.0: scorecards + proposal queue · v2.x: cross-estate governance.",
                  color=ACCENT_TEAL,
                  image_box=(Inches(1.0), Inches(2.6), Inches(11.3), Inches(2.5)))
    # extra bullets under the arrow
    # (rendered into the same slide via re-fetching last slide)

    # 20. Benefits
    slide_benefits(prs)

    # 21. Reference
    slide_reference(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT}")


if __name__ == "__main__":
    build()
